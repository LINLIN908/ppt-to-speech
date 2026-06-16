#!/usr/bin/env python3
"""Write generated speech scripts into a deck's speaker notes.

Usage:
    python write_notes.py input.pptx scripts.json --output input_with_notes.pptx

scripts.json format (slide numbers are 1-based strings or ints):
{
  "1": "Script text for slide 1 ...",
  "2": "Script text for slide 2 ..."
}

Behavior:
- Never modifies the input file; always writes to --output.
- If a slide already has author notes, the generated script is prepended and
  the original notes are preserved below a "---" separator.
- Slides absent from scripts.json are left untouched.

Requires: python-pptx
"""
import argparse
import json
import sys


def main():
    try:
        from pptx import Presentation
    except ImportError:
        sys.exit("python-pptx is required: pip install python-pptx")

    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("scripts_json")
    ap.add_argument("--output", required=True)
    args = ap.parse_args()

    if args.output == args.pptx:
        sys.exit("Refusing to overwrite the input file; choose a different --output.")

    with open(args.scripts_json, encoding="utf-8") as f:
        scripts = {int(k): v for k, v in json.load(f).items()}

    prs = Presentation(args.pptx)
    written = 0
    for i, slide in enumerate(prs.slides, start=1):
        if i not in scripts:
            continue
        notes_slide = slide.notes_slide  # creates one if absent
        tf = notes_slide.notes_text_frame
        existing = tf.text.strip()
        new_text = scripts[i].strip()
        if existing:
            new_text = f"{new_text}\n\n---\n{existing}"
        tf.text = new_text
        written += 1

    prs.save(args.output)
    print(f"Wrote notes for {written} slide(s) -> {args.output}")


if __name__ == "__main__":
    main()
