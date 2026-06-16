#!/usr/bin/env python3
"""Extract per-slide content from a .pptx into JSON.

Usage:
    python extract_slides.py input.pptx --output slides.json

Output format:
{
  "slide_count": 12,
  "slides": [
    {
      "number": 1,
      "title": "...",
      "body_text": ["...", "..."],
      "tables": [[["r1c1","r1c2"],["r2c1","r2c2"]]],
      "image_alt_texts": ["..."],
      "existing_notes": "...",
      "is_text_thin": false
    }
  ]
}

`is_text_thin` flags slides with fewer than ~10 words of extractable text —
the caller should ask the user about those slides rather than invent content.

Requires: python-pptx  (pip install python-pptx)
"""
import argparse
import json
import sys


def extract_text_frames(shape, body):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if text:
                body.append(text)
    if shape.shape_type == 6:  # GROUP
        for sub in shape.shapes:
            extract_text_frames(sub, body)


def main():
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        sys.exit("python-pptx is required: pip install python-pptx")

    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--output", default="slides.json")
    args = ap.parse_args()

    prs = Presentation(args.pptx)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body, tables, alts = [], [], []
        for shape in slide.shapes:
            if shape == slide.shapes.title and shape.has_text_frame:
                title = shape.text_frame.text.strip()
                continue
            if shape.has_table:
                tables.append(
                    [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                )
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                alt = (shape._element.xpath(".//*[@descr]") or [None])[0]
                if alt is not None and alt.get("descr"):
                    alts.append(alt.get("descr"))
                continue
            extract_text_frames(shape, body)

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        word_count = len((title + " " + " ".join(body)).split()) + len(title + "".join(body)) // 4
        slides.append({
            "number": i,
            "title": title,
            "body_text": body,
            "tables": tables,
            "image_alt_texts": alts,
            "existing_notes": notes,
            "is_text_thin": word_count < 10 and not tables,
        })

    out = {"slide_count": len(slides), "slides": slides}
    with open(args.output, "w", encoding="utf-8") as f:
        json.dump(out, f, ensure_ascii=False, indent=2)
    thin = [s["number"] for s in slides if s["is_text_thin"]]
    print(f"Extracted {len(slides)} slides -> {args.output}")
    if thin:
        print(f"Text-thin slides (ask the user about these): {thin}")


if __name__ == "__main__":
    main()
